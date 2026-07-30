"""
Generate a 12-minute PowerPoint presentation:
Fake News Detection with NLP Classification
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
SLIDE_BG  = RGBColor(0x0F, 0x0F, 0x23)   # even darker
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_R  = RGBColor(0xE7, 0x4C, 0x3C)   # red (fake)
ACCENT_G  = RGBColor(0x2E, 0xCC, 0x71)   # green (real)
ACCENT_B  = RGBColor(0x34, 0x98, 0xDB)   # blue
ACCENT_Y  = RGBColor(0xF3, 0x9C, 0x12)   # gold
ACCENT_P  = RGBColor(0x9B, 0x59, 0xB6)   # purple
GRAY      = RGBColor(0xBB, 0xBB, 0xBB)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──
def slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=WHITE, bold_first=False, spacing=Pt(6)):
    """lines is a list of (text, font_size_override, bold_override, color_override) or str"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, fs, bld, clr = line, font_size, False, color
        else:
            txt = line[0]
            fs  = line[1] if len(line) > 1 and line[1] else font_size
            bld = line[2] if len(line) > 2 and line[2] is not None else False
            clr = line[3] if len(line) > 3 and line[3] else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Calibri"
        p.space_after = spacing
    return tf

def add_title_bar(slide, title_text, subtitle_text=None):
    """Dark top bar with title"""
    # Colored accent line at very top
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_B
    line.line.fill.background()

    add_textbox(slide, 0.7, 0.25, 11.5, 0.9, title_text,
                font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, 0.7, 1.0, 11.5, 0.5, subtitle_text,
                    font_size=16, color=GRAY, alignment=PP_ALIGN.LEFT)

def add_page_number(slide, num):
    add_textbox(slide, 12.3, 7.0, 0.8, 0.4, str(num),
                font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

def make_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=ACCENT_B):
    """data is list of lists; first row is header"""
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(13)
                para.font.name = "Calibri"
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                else:
                    para.font.color.rgb = WHITE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x36)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x2A)
    return table

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide_bg(slide)

# Big centered title
add_textbox(slide, 1.0, 1.8, 11.3, 1.5,
            "Fake News Detection\nwith NLP Classification",
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 3.6, 11.3, 0.6,
            "TF-IDF  •  Word2Vec  •  Random Forest  •  XGBoost  •  Gradient Boosting",
            font_size=20, color=ACCENT_B, alignment=PP_ALIGN.CENTER)

# Divider line
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.4), Inches(4.3), Inches(0.04))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT_B; div.line.fill.background()

add_textbox(slide, 1.0, 4.7, 11.3, 0.5,
            "Ironhack Data Analytics Bootcamp  •  Week 7  •  Project 3",
            font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 6.5, 11.3, 0.4,
            "Models: Random Forest 🥇 | XGBoost | Hist Gradient Boosting | Gradient Boosting",
            font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — WHERE NLP CLASSIFICATION SITS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Where NLP Classification Sits in the Landscape", "The 5 Eras of Text Classification")
add_page_number(slide, 2)

# Timeline boxes
eras = [
    ("1. Rule-Based\n1950s–1990s", "Hand-crafted regex\nKeyword matching\nIF 'fake' THEN label=0", ACCENT_R),
    ("2. Classical ML\n2000s–2010s", "TF-IDF + SVMs\nRandom Forests\nNaive Bayes", ACCENT_Y),
    ("3. Embeddings\n2013–2018", "Word2Vec, GloVe\nFastText\n← We are here", ACCENT_G),
    ("4. Deep RNNs\n2014–2019", "LSTMs, GRUs\nBiLSTM + Attention", ACCENT_B),
    ("5. Transformers\n2018–Today 🚀", "BERT, GPT, LLMs\nRoBERTa, T5\nCurrent SOTA", ACCENT_P),
]
for i, (title, desc, color) in enumerate(eras):
    x = 0.5 + i * 2.5
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.3), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x3A)
    box.line.color.rgb = color; box.line.width = Pt(2)
    add_textbox(slide, x + 0.15, 1.9, 2.0, 0.7, title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 2.5, 2.0, 0.8, desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < 4:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.35), Inches(2.45), Inches(0.15), Inches(0.3))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = LIGHT_GRAY; arrow.line.fill.background()

# Highlight box for our project
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(3.7), Inches(5.0), Inches(1.0))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x15, 0x3D, 0x1A)
box.line.color.rgb = ACCENT_G; box.line.width = Pt(2.5)
add_textbox(slide, 5.7, 3.82, 4.6, 0.8,
            "🎯 This Project: Classical ML + Embeddings\n91.2% accuracy • No GPU needed • Fully interpretable",
            font_size=15, color=ACCENT_G, bold=True, alignment=PP_ALIGN.CENTER)

# Why classical still matters
add_multiline(slide, 0.7, 5.2, 11.5, 2.0, [
    ("Why Classical ML Still Dominates in Production:", 18, True, ACCENT_Y),
    ("• Efficiency — 8,000 sparse features vs BERT's 768-dim × 512 tokens per sample", 14, False, GRAY),
    ("• Interpretability — You can explain exactly WHY a headline was classified as fake", 14, False, GRAY),
    ("• Speed — Trains in seconds on CPU; inference in microseconds per headline", 14, False, GRAY),
    ("• Cost — No GPU cluster needed; runs on a $5/month cloud instance", 14, False, GRAY),
], font_size=14, spacing=Pt(4))

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT & DATASET
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Problem Statement & Dataset")
add_page_number(slide, 3)

add_textbox(slide, 0.7, 1.6, 11.5, 0.5,
            'Goal: Build a classifier to distinguish Fake (0) vs Real (1) news headlines',
            font_size=20, color=WHITE, bold=False)

# Dataset format
make_table(slide, 0.7, 2.3, 12.0, 1.5, 5, 3, [
    ["", "Fake News (label = 0)", "Real News (label = 1)"],
    ["Count", "~5,150 (51.5%)", "~4,850 (48.5%)"],
    ["Avg Words", "7.2 words", "8.1 words"],
    ["Example", '"Trump sends embarrassing New Year message"', '"Fed raises interest rates by 0.25%"'],
    ["Tone", "Emotional, clickbait, ALL-CAPS, sensational", "Factual, institutional, measured"],
], col_widths=[1.8, 5.1, 5.1])

# Format notes
add_multiline(slide, 0.7, 4.2, 11.5, 1.5, [
    ("Format:  tab-separated file — <label>\\t<headline>", 16, True, ACCENT_B),
    ("✓ Nearly balanced dataset — no severe class imbalance (we'll still test SMOTE to confirm)", 14, False, GRAY),
    ("✓ No missing values, no duplicate headlines", 14, False, GRAY),
    ("✓ Short text — headlines are 5–15 words, ideal for TF-IDF and Word2Vec", 14, False, GRAY),
], spacing=Pt(6))

# Key insight box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.5), Inches(11.8), Inches(1.5))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x35)
box.line.color.rgb = ACCENT_Y; box.line.width = Pt(1.5)
add_multiline(slide, 1.0, 5.65, 11.2, 1.2, [
    ("💡 Key Challenge:", 16, True, ACCENT_Y),
    ("Fake and real headlines share vocabulary — both discuss politics, Trump, White House...", 14, False, GRAY),
    ('Classification must rely on subtle linguistic patterns, not obvious keywords.', 14, False, GRAY),
    ('This is why simple keyword-matching (Era 1) fails — we need statistical feature extraction.', 14, False, GRAY),
], spacing=Pt(3))

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — TEXT PREPROCESSING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Text Preprocessing Pipeline", "Every decision here impacts final accuracy — this is the critical step")
add_page_number(slide, 4)

# Pipeline steps as connected boxes
steps = [
    ("1. Lowercase", '"Trump" → "trump"', ACCENT_R),
    ("2. Remove\nURLs/HTML", "regex: http...\n<html>", ACCENT_Y),
    ("3. Remove\nPunctuation", "[^a-zA-Z\\s]", ACCENT_G),
    ("4. Remove\nNumbers", "r\"\\d+\"", ACCENT_B),
    ("5. Tokenize", 'nltk.word_tokenize', ACCENT_P),
    ("6. Remove\nStopwords", "the, a, is, said...", ACCENT_R),
    ("7. Lemmatize", "running → run", ACCENT_Y),
    ("8. Cleaned\nText", "Ready for features!", ACCENT_G),
]
for i, (title, desc, color) in enumerate(steps):
    x = 0.4 + i * 1.6
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.7), Inches(1.45), Inches(1.3))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x3A)
    box.line.color.rgb = color; box.line.width = Pt(1.5)
    add_textbox(slide, x + 0.05, 1.78, 1.35, 0.55, title, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.05, 2.3, 1.35, 0.6, desc, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    if i < 7:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.45), Inches(2.2), Inches(0.15), Inches(0.25))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = LIGHT_GRAY; arrow.line.fill.background()

# Technical decisions table
make_table(slide, 0.7, 3.4, 12.0, 2.6, 5, 3, [
    ["Step", "Technique & Code", "Why This Matters"],
    ["Stopwords", "NLTK + custom: {'said', 'says', 'new', 'us', 'would'}", "News-specific: 'said' appears everywhere, adds zero signal"],
    ["Lemmatization", "WordNetLemmatizer — reduces to base form", "Better than stemming for news: 'running'→'run' not 'runn'"],
    ["Min Word Length", "Keep words with len > 2", "Filters noise: 'a', 'an', 'at', 'be', 'go'"],
    ["Empty Handling", "Fill empty texts with 'empty' placeholder", "Prevents NaN errors in vectorization"],
], col_widths=[2.2, 5.3, 4.5])

add_multiline(slide, 0.7, 6.3, 11.5, 0.9, [
    ("Before:  \"Donald Trump SENDS OUT embarrassing New Year's Eve message — THIS is disturbing!!!\"", 12, False, LIGHT_GRAY),
    ("After:    \"donald trump send embarrassing new year eve message disturbing\"", 12, True, ACCENT_G),
], spacing=Pt(2))

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — EDA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Exploratory Data Analysis", "Word clouds, N-grams, and text length distributions")
add_page_number(slide, 5)

# Word cloud note
add_textbox(slide, 0.7, 1.5, 5.5, 0.4, "🔴 Fake News Top Words", font_size=18, color=ACCENT_R, bold=True)
add_textbox(slide, 7.0, 1.5, 5.5, 0.4, "🟢 Real News Top Words", font_size=18, color=ACCENT_G, bold=True)

add_textbox(slide, 0.7, 1.9, 5.5, 0.8,
            "trump • video • watch • obama • police\nblack • racist • donald • president • like",
            font_size=14, color=GRAY)

add_textbox(slide, 7.0, 1.9, 5.5, 0.8,
            "trump • north • korea • house • china\npresident • white • new • state • says",
            font_size=14, color=GRAY)

# N-gram findings
add_multiline(slide, 0.7, 3.0, 11.5, 1.5, [
    ("🔍 Key EDA Findings:", 20, True, ACCENT_Y),
    ("• Fake & Real share vocabulary heavily (both discuss Trump, White House, politics)", 15, False, GRAY),
    ("• Fake headlines use emotional/clickbait: 'watch', 'video', 'breaking', 'this is disturbing'", 15, False, GRAY),
    ("• Real headlines use institutional terms: 'senate', 'republican', 'administration'", 15, False, GRAY),
    ("• Bigram 'white house' appears in BOTH — no single word is a reliable signal", 15, False, GRAY),
], spacing=Pt(5))

# Word count comparison
make_table(slide, 0.7, 4.8, 12.0, 1.2, 4, 3, [
    ["Metric", "Fake News", "Real News"],
    ["Avg Word Count (cleaned)", "~5.2 words", "~5.8 words"],
    ["Avg Char Count (original)", "~45 chars", "~52 chars"],
    ["Top Bigrams", "white house, north korea, donald trump", "white house, united states, north korea"],
], col_widths=[3.0, 4.5, 4.5])

add_multiline(slide, 0.7, 6.3, 11.5, 0.9, [
    ("💡 Implication: Classification must learn subtle distributional differences, not obvious keywords.", 15, True, ACCENT_Y),
], spacing=Pt(3))

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — FEATURE EXTRACTION: TF-IDF
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Feature Extraction — TF-IDF (Bag-of-Words)", "Term Frequency × Inverse Document Frequency with N-grams")
add_page_number(slide, 6)

# Formula
add_textbox(slide, 0.7, 1.5, 11.5, 0.7,
            "TF-IDF(t, d) = TF(t,d) × log(N / DF(t))",
            font_size=24, color=ACCENT_B, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.7, 2.1, 11.5, 0.4,
            "Term Frequency — how often word t appears in doc d    ×    Inverse Document Frequency — how rare is t across all docs",
            font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Configuration
add_multiline(slide, 0.7, 2.7, 11.5, 1.2, [
    ("Configuration:", 18, True, WHITE),
    ("ngram_range=(1,3)  →  Unigrams + Bigrams + Trigrams  —  'white house' ≠ 'white' + 'house'", 15, False, ACCENT_B),
    ("max_features=8000  →  Limit vocabulary to top 8,000 features to control dimensionality", 15, False, ACCENT_B),
    ("min_df=3, max_df=0.85  →  Ignore very rare words (<3 docs) and corpus-wide terms (>85% docs)", 15, False, ACCENT_B),
    ("sublinear_tf=True  →  1 + log(tf) — dampens effect of high-frequency words", 15, False, ACCENT_B),
], spacing=Pt(3))

# Why n-grams
add_multiline(slide, 0.7, 4.2, 5.5, 2.8, [
    ("Why N-grams (1,3)?", 18, True, ACCENT_Y),
    ("", 6, False, GRAY),
    ('Without n-grams:', 14, True, GRAY),
    ('  "white" + "house" → two separate features', 13, False, GRAY),
    ('  "north" + "korea" → no connection', 13, False, GRAY),
    ("", 6, False, GRAY),
    ('With n-grams:', 14, True, ACCENT_G),
    ('  "white house" → one feature (The White House)', 13, False, ACCENT_G),
    ('  "north korea" → one feature (the country)', 13, False, ACCENT_G),
    ('  "donald trump jr" → one feature (specific person)', 13, False, ACCENT_G),
], spacing=Pt(2))

# Result matrix
add_multiline(slide, 7.0, 4.2, 5.5, 2.8, [
    ("Result:", 18, True, WHITE),
    ("", 6, False, GRAY),
    ("Training matrix: 8,000+ rows × 8,000 cols", 14, False, GRAY),
    ("Vocabulary size: 8,000 terms/ngrams", 14, False, GRAY),
    ("Sparsity: ~0.12% non-zero (highly sparse)", 14, False, GRAY),
    ("", 8, False, GRAY),
    ("📊 Each row = one headline vectorized", 14, False, LIGHT_GRAY),
    ("📊 Each column = one TF-IDF weighted term", 14, False, LIGHT_GRAY),
    ("📊 Most values are 0.0 (sparse matrix)", 14, False, LIGHT_GRAY),
], spacing=Pt(2))

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — FEATURE EXTRACTION: Word2Vec
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Feature Extraction — Word2Vec Embeddings (Comparison)", "Dense semantic vectors trained on the headline corpus itself")
add_page_number(slide, 7)

# Word2Vec config
add_multiline(slide, 0.7, 1.5, 6.0, 2.0, [
    ("Word2Vec Configuration:", 18, True, WHITE),
    ("vector_size=200    →  200-dimensional dense vectors", 15, False, ACCENT_G),
    ("window=5           →  Context window of 5 words around target", 15, False, ACCENT_G),
    ("sg=1 (Skip-gram)   →  Better than CBOW for small datasets", 15, False, ACCENT_G),
    ("min_count=2        →  Ignore words appearing only once", 15, False, ACCENT_G),
    ("epochs=20          →  Multiple passes over the corpus", 15, False, ACCENT_G),
], spacing=Pt(3))

# Document vector formula
add_textbox(slide, 7.5, 1.5, 5.5, 0.4,
            "Document Vector:", font_size=18, color=WHITE, bold=True)
add_textbox(slide, 7.5, 1.9, 5.5, 0.8,
            "doc_vec = average( word_vec(w₁),\n                     word_vec(w₂),\n                     ... )",
            font_size=15, color=ACCENT_B, bold=True)

add_textbox(slide, 7.5, 2.8, 5.5, 0.5,
            "Similar words to 'trump':\n  obama: 0.82  |  donald: 0.79\n  president: 0.74  |  clinton: 0.71",
            font_size=13, color=GRAY)

# Comparison table
make_table(slide, 0.7, 3.8, 12.0, 2.4, 6, 3, [
    ["Aspect", "TF-IDF", "Word2Vec"],
    ["Dimensionality", "8,000 features (sparse)", "200 features (dense) ✨"],
    ["Semantics", "❌ No — 'king' ≠ 'monarch'", "✅ Yes — captures similarity"],
    ["Interpretability", "✅ High — see exact words", "❌ Low — abstract vectors"],
    ["Training Time", "Seconds", "~30 seconds"],
    ["Accuracy (Random Forest)", "91.26% 🥇", "91.16% 🥈"],
], col_widths=[2.8, 4.6, 4.6])

add_multiline(slide, 0.7, 6.5, 11.5, 0.8, [
    ("🔥 Word2Vec matches TF-IDF accuracy with 40× fewer features! Ideal for low-latency production.", 16, True, ACCENT_Y),
], spacing=Pt(2))

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — MODEL TRAINING STRATEGY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Model Training Strategy — 3-Phase Approach", "Baselines → Hyperparameter Tuning → SMOTE Balancing Test")
add_page_number(slide, 8)

# Phase 1
add_textbox(slide, 0.7, 1.5, 3.8, 0.4, "PHASE 1: Baselines", font_size=18, color=ACCENT_B, bold=True)
make_table(slide, 0.7, 1.95, 3.8, 2.0, 5, 3, [
    ["Model", "Accuracy", "F1"],
    ["Random Forest", "~91.2%", "0.912"],
    ["XGBoost", "~88.0%", "0.880"],
    ["Hist GB", "~87.0%", "0.870"],
    ["Gradient B.", "~79.7%", "0.797"],
], col_widths=[1.3, 1.2, 1.3], header_color=ACCENT_R)

# Phase 2
add_textbox(slide, 5.0, 1.5, 4.0, 0.4, "PHASE 2: Hyperparameter Tuning", font_size=18, color=ACCENT_Y, bold=True)
add_multiline(slide, 5.0, 1.95, 4.0, 2.0, [
    ("Method: RandomizedSearchCV", 13, True, WHITE),
    ("• n_iter=5 random combinations", 12, False, GRAY),
    ("• cv=StratifiedKFold(2)", 12, False, GRAY),
    ("• scoring='f1_macro'", 12, False, GRAY),
    ("", 4, False, GRAY),
    ("Why F1-macro?", 13, True, ACCENT_Y),
    ("• Penalizes class favoritism", 12, False, GRAY),
    ("• Balances precision & recall", 12, False, GRAY),
    ("• Critical when both FP and FN matter", 12, False, GRAY),
], spacing=Pt(1))

# Phase 3
add_textbox(slide, 9.5, 1.5, 3.5, 0.4, "PHASE 3: SMOTE", font_size=18, color=ACCENT_G, bold=True)
add_multiline(slide, 9.5, 1.95, 3.5, 2.0, [
    ("SMOTE = Synthetic Minority", 13, True, WHITE),
    ("Oversampling Technique", 13, True, WHITE),
    ("", 4, False, GRAY),
    ("Creates synthetic samples:", 12, False, GRAY),
    ("x_new = x_i + λ·(x_neighbor − x_i)", 11, False, ACCENT_B),
    ("", 4, False, GRAY),
    ("Result: ❌ Did NOT improve", 13, True, ACCENT_R),
    ("Data already 51.5/48.5 balanced!", 12, False, GRAY),
    ("Lesson: Always verify imbalance", 12, False, GRAY),
], spacing=Pt(1))

# Tuned params table
add_textbox(slide, 0.7, 4.3, 11.5, 0.4, "Hyperparameters Tuned per Model:", font_size=16, color=WHITE, bold=True)
make_table(slide, 0.7, 4.7, 12.0, 1.8, 5, 2, [
    ["Model", "Parameters Tuned"],
    ["Random Forest", "n_estimators: [100,200,300], max_depth: [None,20,50], min_samples_split: [2,10], min_samples_leaf: [1,4], max_features: ['sqrt','log2']"],
    ["XGBoost", "n_estimators: [100,200,300], max_depth: [3,5,9], learning_rate: [0.01,0.1,0.2], subsample: [0.6,0.8,1.0], gamma: [0,0.1]"],
    ["Hist GB", "max_iter: [100,200,300], learning_rate: [0.01,0.1,0.2], max_depth: [None,10,50], l2_regularization: [0,0.5]"],
    ["Gradient B.", "n_estimators: [100,200,300], learning_rate: [0.01,0.1,0.2], max_depth: [3,5,7], subsample: [0.6,1.0]"],
], col_widths=[2.2, 9.8])

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — FINAL RESULTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Final Results — Model Comparison", "TF-IDF vs Word2Vec • Baseline vs Tuned • SMOTE impact")
add_page_number(slide, 9)

make_table(slide, 0.5, 1.5, 12.3, 2.0, 6, 5, [
    ["Rank", "Model", "Features", "Accuracy", "F1-Macro"],
    ["🥇", "TF-IDF + Random Forest (Tuned)", "8,000", "91.26%", "0.9126"],
    ["🥈", "Word2Vec + Random Forest", "200 ✨", "91.16%", "0.9114"],
    ["🥉", "TF-IDF + Hist Gradient Boosting", "8,000", "90.91%", "0.9091"],
    ["4", "TF-IDF + Gradient Boosting", "8,000", "90.73%", "0.9073"],
    ["5", "TF-IDF + XGBoost", "8,000", "89.28%", "0.8928"],
], col_widths=[0.8, 4.4, 1.8, 1.8, 1.8])

# Confusion matrix interpretation
add_textbox(slide, 0.5, 3.8, 5.8, 0.4, "Confusion Matrix — Best Model (RF):", font_size=18, color=WHITE, bold=True)

make_table(slide, 0.5, 4.2, 5.8, 1.2, 3, 3, [
    ["", "Predicted Fake", "Predicted Real"],
    ["Actual Fake", "89.2% ✅", "10.8% (False Positive)"],
    ["Actual Real", "6.7% (False Negative)", "93.3% ✅"],
], col_widths=[2.0, 1.9, 1.9], header_color=ACCENT_B)

# ROC-AUC
add_multiline(slide, 7.0, 3.8, 5.5, 2.0, [
    ("ROC-AUC: All Models > 0.96", 18, True, ACCENT_G),
    ("• Excellent class separation across all models", 14, False, GRAY),
    ("• Random Forest AUC: 0.971 ± 0.003", 14, False, GRAY),
    ("• Gradient Boosting improved most from tuning:", 14, False, GRAY),
    ("  Baseline F1: 0.797 → Tuned F1: 0.907 (+11%) 🔥", 14, True, ACCENT_Y),
    ("• SMOTE: No gain (data already balanced)", 14, False, GRAY),
    ("• Feature Selection: Best at median threshold (5,200 features)", 14, False, GRAY),
], spacing=Pt(3))

# GB improvement highlight
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(6.0), Inches(5.5), Inches(1.0))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x2A, 0x1A, 0x0A)
box.line.color.rgb = ACCENT_Y; box.line.width = Pt(2)
add_textbox(slide, 7.2, 6.1, 5.1, 0.8,
            "🔥 Biggest Tuner: Gradient Boosting\nBaseline: 79.7% F1  →  Tuned: 90.7% F1\nHyperparameter tuning is NOT optional!",
            font_size=14, color=ACCENT_Y, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — FEATURE IMPORTANCE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Feature Importance — What Makes a Headline 'Fake'?", "Top TF-IDF features driving the Random Forest classifier")
add_page_number(slide, 10)

add_textbox(slide, 0.7, 1.5, 5.5, 0.4, "Top Features → FAKE (0)", font_size=20, color=ACCENT_R, bold=True)
add_textbox(slide, 7.0, 1.5, 5.5, 0.4, "Top Features → REAL (1)", font_size=20, color=ACCENT_G, bold=True)

make_table(slide, 0.7, 2.0, 5.5, 3.0, 8, 2, [
    ["Feature", "Importance"],
    ["trump", "0.023"],
    ["video", "0.018"],
    ["watch", "0.014"],
    ["police", "0.012"],
    ["breaking", "0.009"],
    ["black", "0.008"],
    ["racist", "0.007"],
], col_widths=[2.5, 3.0], header_color=ACCENT_R)

make_table(slide, 7.0, 2.0, 5.5, 3.0, 8, 2, [
    ["Feature", "Importance"],
    ["white house", "0.008"],
    ["senate", "0.007"],
    ["republican", "0.007"],
    ["administration", "0.006"],
    ["congress", "0.006"],
    ["election", "0.006"],
    ["federal", "0.005"],
], col_widths=[2.5, 3.0], header_color=ACCENT_G)

# Pattern analysis
add_multiline(slide, 0.7, 5.3, 11.5, 2.0, [
    ("📊 Pattern Analysis:", 18, True, ACCENT_Y),
    ("• Fake news:  Emotional/clickbait language — 'watch', 'video', 'breaking', 'racist', 'disturbing'", 14, False, GRAY),
    ("• Real news:  Institutional/political terms — 'senate', 'congress', 'administration', 'federal'", 14, False, GRAY),
    ("• Both use 'trump' heavily — the signal is in the co-occurring words, not 'trump' alone", 14, False, GRAY),
    ("• N-grams are critical: 'white house' (bigram) is a REAL signal; 'white' alone means nothing", 14, False, GRAY),
], spacing=Pt(4))

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — KEY TECHNICAL TAKEAWAYS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Key Technical Takeaways", "What we learned from building this classifier")
add_page_number(slide, 11)

takeaways = [
    ("1", "TF-IDF + Random Forest is surprisingly strong", "91.26% accuracy on short-text classification — beats gradient boosting on this dataset", ACCENT_R),
    ("2", "Word2Vec = TF-IDF with 40× fewer features", "Same accuracy (91.2%) with 200 dimensions vs 8,000 — ideal for production at scale", ACCENT_G),
    ("3", "N-grams capture critical phrases", "'white house', 'north korea', 'donald trump jr' — unigrams alone would miss these", ACCENT_B),
    ("4", "SMOTE didn't help — always verify first", "Data was already balanced (51.5/48.5). Blindly applying techniques wastes time.", ACCENT_P),
    ("5", "Hyperparameter tuning is transformative", "Gradient Boosting jumped from 79.7% → 90.7% F1 — never skip tuning!", ACCENT_Y),
    ("6", "Preprocessing matters more than model choice", "Custom stopwords (said, says, new) alone can shift accuracy by 1-2 percentage points", ACCENT_R),
    ("7", "Where would Transformers take us?", "BERT/RoBERTa would likely hit 94–96% — but at 100× the compute cost and lost interpretability", ACCENT_G),
]
for i, (num, title, desc, color) in enumerate(takeaways):
    y = 1.6 + i * 0.82
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.45), Inches(0.45))
    circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
    add_textbox(slide, 0.7, y + 0.05, 0.45, 0.4, num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.3, y - 0.02, 11.0, 0.35, title, font_size=17, color=color, bold=True)
    add_textbox(slide, 1.3, y + 0.32, 11.0, 0.35, desc, font_size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — PRODUCTION PATH
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_title_bar(slide, "Production Path — If This Were Real", "Architecture and next steps for a deployable system")
add_page_number(slide, 12)

# Flow diagram as boxes
flow_steps = [
    ("News\nHeadline", ACCENT_R),
    ("Preprocessing\nPipeline", ACCENT_Y),
    ("TF-IDF\nTransform", ACCENT_G),
    ("Random Forest\nPredict", ACCENT_B),
    ("Score > 0.5?", ACCENT_P),
    ("FAKE", ACCENT_R),
    ("REAL", ACCENT_G),
]
for i, (label, color) in enumerate(flow_steps):
    x = 0.5 + i * 1.8
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(1.6), Inches(1.0))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x3A)
    box.line.color.rgb = color; box.line.width = Pt(2)
    add_textbox(slide, x + 0.05, 1.72, 1.5, 0.8, label, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 4:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.6), Inches(1.95), Inches(0.2), Inches(0.3))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = LIGHT_GRAY; arrow.line.fill.background()

# Confidence routing
add_textbox(slide, 9.6, 1.6, 1.5, 0.5, "→ FAKE\n→ REAL", font_size=14, color=ACCENT_Y, bold=True)

# Human review path
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(2.8), Inches(3.2), Inches(0.7))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x2A, 0x1A, 0x0A)
box.line.color.rgb = ACCENT_Y; box.line.width = Pt(2)
add_textbox(slide, 6.35, 2.88, 2.9, 0.5,
            "⚠️  Confidence < 0.7?\n→ Human Review Queue",
            font_size=13, color=ACCENT_Y, bold=True, alignment=PP_ALIGN.CENTER)

# Next steps
add_multiline(slide, 0.7, 3.9, 5.5, 3.2, [
    ("Next Steps for Improvement:", 18, True, WHITE),
    ("", 4, False, GRAY),
    ("1. Fine-tune distilBERT", 15, True, ACCENT_B),
    ("   3× smaller than BERT, ~94% expected", 13, False, GRAY),
    ("", 4, False, GRAY),
    ("2. Ensemble Voting Classifier", 15, True, ACCENT_B),
    ("   RF + XGBoost + HGB → reduce variance", 13, False, GRAY),
    ("", 4, False, GRAY),
    ("3. Add metadata features", 15, True, ACCENT_B),
    ("   Source domain, publish time, author history", 13, False, GRAY),
    ("", 4, False, GRAY),
    ("4. Deploy as FastAPI microservice", 15, True, ACCENT_B),
    ("   < 1ms inference per headline on CPU", 13, False, GRAY),
], spacing=Pt(1))

add_multiline(slide, 7.0, 3.9, 5.5, 3.2, [
    ("Performance Summary:", 18, True, WHITE),
    ("", 6, False, GRAY),
    ("Model:  Random Forest", 16, True, ACCENT_G),
    ("Features:  TF-IDF (8k n-grams 1-3)", 16, False, GRAY),
    ("Accuracy:  91.26%", 16, True, ACCENT_G),
    ("F1-Macro:  0.9126", 16, False, GRAY),
    ("ROC-AUC:  0.971", 16, False, GRAY),
    ("", 8, False, GRAY),
    ("Predictions saved to:", 14, False, LIGHT_GRAY),
    ("dataset/predictions.csv ✅", 14, True, ACCENT_G),
], spacing=Pt(2))

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)

add_textbox(slide, 1.0, 1.5, 11.3, 1.0,
            "Conclusion",
            font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(2.5), Inches(3.3), Inches(0.04))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT_B; div.line.fill.background()

add_multiline(slide, 1.0, 3.0, 11.3, 3.0, [
    ("NLP Classification with classical ML is alive and well.", 24, True, ACCENT_G),
    ("", 10, False, GRAY),
    ("For short-text tasks like headline classification,", 18, False, GRAY),
    ("TF-IDF + tree-based models deliver 91%+ accuracy —", 18, False, GRAY),
    ("interpretable, fast, and production-ready —", 18, False, GRAY),
    ("without the overhead of deep learning.", 18, False, GRAY),
    ("", 10, False, GRAY),
    ("Word2Vec proves you can match that accuracy", 18, False, GRAY),
    ("with just 200 dimensions. Efficiency matters.", 18, False, GRAY),
], spacing=Pt(4))

add_textbox(slide, 1.0, 6.2, 11.3, 0.6,
            "Final Model: Random Forest + TF-IDF (n-grams 1-3, 8k features)  •  Accuracy: 91.26%  •  F1: 0.9126  •  ROC-AUC: 0.97+",
            font_size=15, color=ACCENT_B, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 6.9, 11.3, 0.4,
            "Thank you!  —  Questions?",
            font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = r"c:\Users\cici\Documents\Ironhack\week 7\D2\project-3-nlp\Fake_News_Detection_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   {len(prs.slides)} slides generated")
